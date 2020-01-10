from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
import glob

prs = Presentation("D:\\Text\\8d\\8D_Classic\\8D_Classic\\65SD BA-Fall time.pptx.pptx")
for slide in prs.slides:
    print(slide)
    print("slide id: {}".format(slide.slide_id))
    for shape in slide.shapes:
        print(shape)
        if isinstance(shape, GraphicFrame):
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        print(cell.text)
        if hasattr(shape, "text"):
            print(shape.text)
        